"""Tests for Document Designer integration in file_tools.

Covers:

1. **In-process sandbox** — `run_in_sandbox(..., force_in_process=True)`
   executes validated code and returns its bytes. This is the path the
   PyInstaller-frozen desktop sidecar uses (sys.executable is the
   bundle, not a real Python).
2. **Designer routing** — `_try_designer` invokes PptxDesignerSkill /
   DocxDesignerSkill when format is pptx/docx and an LLM is on the
   agent. Falls back gracefully on every failure mode.
3. **Pending-op cache** — when Designer succeeds at queue time, the
   bytes are stored as base64 on the pending op and apply_pending_op
   writes them verbatim (no second LLM round-trip).
4. **Apply-time fallback** — when no designer_bytes_b64 is present,
   apply_pending_op falls back to export_as.
"""
from __future__ import annotations

import base64
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from praxia.agent.file_tools import (
    _try_designer,
    apply_pending_op,
    workspace_tools,
)
from praxia.core.llm import LLMResponse


# Minimal python-pptx snippet that the AST validator + in-process
# runner can execute end-to-end. We stub the Designer's LLM to return
# this verbatim so the test doesn't require a real python-pptx install
# unless the env has one; it's gated by importorskip below.
_PPTX_CODE = """
import io
from pptx import Presentation

def build():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test deck"
    buf = io.BytesIO()
    prs.save(buf)
    _emit(buf.getvalue())

build()
"""


# ---------------------------------------------------------------------------
# In-process sandbox
# ---------------------------------------------------------------------------


class TestInProcessSandbox:
    def test_force_in_process_round_trip(self):
        pytest.importorskip("pptx")
        from praxia.skills.document_designer.sandbox import run_in_sandbox
        result = run_in_sandbox(_PPTX_CODE, force_in_process=True, timeout_s=10.0)
        assert result.bytes
        # PPTX files are zipfiles → start with PK header.
        assert result.bytes[:2] == b"PK"

    def test_in_process_still_validates_ast(self):
        from praxia.skills.document_designer.sandbox import (
            SandboxValidationError,
            run_in_sandbox,
        )
        with pytest.raises(SandboxValidationError):
            run_in_sandbox("import os\nos.system('rm -rf /')\n_emit(b'x')", force_in_process=True)

    def test_in_process_raises_when_no_emit(self):
        from praxia.skills.document_designer.sandbox import (
            SandboxError,
            run_in_sandbox,
        )
        with pytest.raises(SandboxError):
            # Validates fine; runs fine; just never calls _emit().
            run_in_sandbox("x = 1 + 1\n", force_in_process=True)


# ---------------------------------------------------------------------------
# Designer routing — _try_designer happy & sad paths
# ---------------------------------------------------------------------------


def _stub_llm_returning(code: str):
    """A fake LLM that returns `code` on every complete() call."""
    llm = MagicMock()
    llm.complete.return_value = LLMResponse(
        text=code, model="stub", usage={}, raw={}, tool_calls=[]
    )
    return llm


class TestTryDesigner:
    def test_returns_bytes_when_codegen_succeeds(self):
        pytest.importorskip("pptx")
        agent = MagicMock()
        agent.llm = _stub_llm_returning(_PPTX_CODE)
        # Force in-process so the test doesn't depend on a subprocess
        # Python being on PATH.
        with patch(
            "praxia.skills.document_designer.codegen.run_in_sandbox",
            side_effect=lambda code, **kw: __import__(
                "praxia.skills.document_designer.sandbox", fromlist=["run_in_sandbox"]
            ).run_in_sandbox(code, force_in_process=True, **{k: v for k, v in kw.items() if k != "force_in_process"}),
        ):
            data, err = _try_designer(agent, "pptx", "Brief: a one-slide test deck.")
        assert err is None, err
        assert data and data[:2] == b"PK"

    def test_returns_error_when_no_llm(self):
        agent = MagicMock()
        agent.llm = None
        data, err = _try_designer(agent, "pptx", "anything")
        assert data is None
        assert "no llm" in (err or "")

    def test_returns_error_on_empty_brief(self):
        agent = MagicMock()
        agent.llm = _stub_llm_returning(_PPTX_CODE)
        data, err = _try_designer(agent, "pptx", "   ")
        assert data is None
        assert "empty" in (err or "")

    def test_returns_error_when_designer_raises(self):
        # Codegen exhausts retries → designer raises RuntimeError, which
        # we want to swallow into a (None, error_str) tuple so the
        # caller can fall back to the plain exporter.
        agent = MagicMock()
        agent.llm = _stub_llm_returning("import os\nos.system('boom')\n_emit(b'x')")  # AST will reject
        data, err = _try_designer(agent, "pptx", "deck about Q4 revenue")
        assert data is None
        assert err is not None


# ---------------------------------------------------------------------------
# Pending-op cache → apply_pending_op uses pre-rendered bytes
# ---------------------------------------------------------------------------


class TestApplyUsesDesignerBytes:
    def test_apply_writes_designer_bytes_verbatim(self, tmp_path: Path):
        # Synthetic bytes — apply doesn't care that this isn't a real
        # pptx; it just writes whatever's in designer_bytes_b64.
        fake_bytes = b"PK\x03\x04 not really a deck but bytes are bytes"
        op = {
            "op": "render_document",
            "path": "out/deck.pptx",
            "format": "pptx",
            "source_markdown": "# Title\n\n- bullet",
            "designer_bytes_b64": base64.b64encode(fake_bytes).decode("ascii"),
            "designer_used": True,
        }
        res = apply_pending_op(tmp_path, op)
        assert res.get("applied") is True
        written = (tmp_path / "out" / "deck.pptx").read_bytes()
        assert written == fake_bytes

    def test_apply_falls_back_when_no_designer_bytes(self, tmp_path: Path):
        # No designer_bytes_b64 → apply falls back to export_as. For an
        # md format the basic exporter is just a passthrough, so we
        # can assert on the written content without spinning up
        # python-pptx.
        op = {
            "op": "render_document",
            "path": "out/notes.md",
            "format": "md",
            "source_markdown": "# Title\n\nbody",
        }
        res = apply_pending_op(tmp_path, op)
        assert res.get("applied") is True
        # Markdown exporter is a passthrough — content preserved.
        written = (tmp_path / "out" / "notes.md").read_text(encoding="utf-8")
        assert "Title" in written and "body" in written


# ---------------------------------------------------------------------------
# End-to-end: render_document tool queues a designer-rendered op
# ---------------------------------------------------------------------------


class TestRenderDocumentToolQueuesDesignerBytes:
    def test_pptx_render_tool_caches_designer_bytes_on_queue(self, tmp_path: Path):
        pytest.importorskip("pptx")
        agent = MagicMock()
        agent.llm = _stub_llm_returning(_PPTX_CODE)
        sink: list[dict] = []
        tools = workspace_tools(
            tmp_path,
            require_confirmation=True,
            pending_sink=sink,
        )
        # Force in-process sandbox so the test doesn't shell out.
        with patch(
            "praxia.skills.document_designer.codegen.run_in_sandbox",
            side_effect=lambda code, **kw: __import__(
                "praxia.skills.document_designer.sandbox", fromlist=["run_in_sandbox"]
            ).run_in_sandbox(code, force_in_process=True, **{k: v for k, v in kw.items() if k != "force_in_process"}),
        ):
            result = tools["render_document"].handler(
                agent,
                path="deck.pptx",
                format="pptx",
                source_markdown="Brief: a one-slide deck about widget sales.",
            )
        assert result.get("pending") is True
        assert sink, "queue should have one op"
        op = sink[0]
        assert op["op"] == "render_document"
        assert op["format"] == "pptx"
        assert op.get("designer_used") is True
        assert "designer_bytes_b64" in op
        # And applying that op writes the cached bytes — same content
        # the designer produced, no re-render.
        apply_res = apply_pending_op(tmp_path, op)
        assert apply_res.get("applied") is True
        on_disk = (tmp_path / "deck.pptx").read_bytes()
        assert on_disk[:2] == b"PK"  # valid pptx header
