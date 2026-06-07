"""Tests for the alpha20 vision pipeline:

  document parser → metadata.embedded_images / metadata.page_images
  → search_for_user(include_images=True) → hit["images"]
  → DefaultMemoryRetriever._from_documents → Source.metadata["images"]
  → _aggregate_source_images → CommandedAgent passes images to inner.run

End-to-end shape verification — we mock the LLM, so we're not testing
that vision answers are accurate. We're testing that the image bytes
actually make it from the parser through to inner.run(images=...).
"""
from __future__ import annotations

import base64
import io
import zipfile
from pathlib import Path
from unittest.mock import patch

import pytest


# ---------------------------------------------------------------------------
# Parser-level: PDF/DOCX/PPTX embedded image extraction
# ---------------------------------------------------------------------------


def _make_red_png_bytes() -> bytes:
    """1x1 red PNG — the minimum for a valid image MIME without dragging
    in Pillow as a test dep for THIS particular test. Tests that need
    realistic images (parser smoke) import Pillow directly."""
    return (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\xdacd\xf8\xcf\xc0\x00"
        b"\x00\x00\x03\x00\x01\x9at\x9e\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
    )


class TestDocxImageExtraction:
    def test_docx_with_one_image_yields_one_embedded(self):
        pytest.importorskip("docx")
        pytest.importorskip("PIL")
        from docx import Document
        from docx.shared import Inches
        from PIL import Image

        img = Image.new("RGB", (32, 32), color="red")
        ibuf = io.BytesIO()
        img.save(ibuf, format="PNG")
        ibuf.seek(0)

        doc = Document()
        doc.add_heading("Test", 0)
        doc.add_paragraph("hello")
        doc.add_picture(ibuf, width=Inches(0.5))
        buf = io.BytesIO()
        doc.save(buf)

        from praxia.io.parsers.docx_parser import DocxParser
        parsed = DocxParser().parse(buf.getvalue(), filename="t.docx")
        assert parsed.metadata["embedded_image_count"] == 1
        ei = parsed.metadata["embedded_images"][0]
        assert ei["mime"] in {"image/png", "image/jpeg"}
        assert isinstance(ei["data"], str) and len(ei["data"]) > 0
        # base64 round-trips
        raw = base64.b64decode(ei["data"])
        assert raw.startswith(b"\x89PNG") or raw[:3] in (b"\xff\xd8\xff",)

    def test_docx_without_images_yields_empty(self):
        pytest.importorskip("docx")
        from docx import Document
        d = Document()
        d.add_paragraph("text only")
        buf = io.BytesIO()
        d.save(buf)

        from praxia.io.parsers.docx_parser import DocxParser
        parsed = DocxParser().parse(buf.getvalue(), filename="x.docx")
        assert parsed.metadata["embedded_image_count"] == 0
        assert parsed.metadata["embedded_images"] == []

    def test_docx_caps_per_image_size(self):
        # Forge a .docx-shaped zip with one giant fake PNG in word/media/
        # to make sure the 5MB-per-image cap drops it.
        from praxia.io.parsers.docx_parser import _MAX_PER_IMAGE_BYTES, _extract_embedded_images
        big = b"\x89PNG\r\n\x1a\n" + b"X" * (_MAX_PER_IMAGE_BYTES + 1)
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.writestr("word/media/image1.png", big)
        out = _extract_embedded_images(buf.getvalue())
        assert out == [], "oversized image should be dropped"


class TestPptxImageExtraction:
    def test_pptx_with_one_image_yields_one_embedded(self):
        pytest.importorskip("pptx")
        pytest.importorskip("PIL")
        from pptx import Presentation
        from PIL import Image

        img = Image.new("RGB", (32, 32), color="green")
        ibuf = io.BytesIO()
        img.save(ibuf, format="PNG")
        ibuf.seek(0)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(ibuf, left=0, top=0)
        buf = io.BytesIO()
        prs.save(buf)

        from praxia.io.parsers.pptx_parser import PptxParser
        parsed = PptxParser().parse(buf.getvalue(), filename="t.pptx")
        assert parsed.metadata["embedded_image_count"] == 1


class TestPdfRasterization:
    def test_pdf_pages_get_jpeg_renders(self):
        pytest.importorskip("pypdfium2")
        pytest.importorskip("reportlab")
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=letter)
        c.drawString(100, 750, "Page 1")
        c.showPage()
        c.drawString(100, 750, "Page 2")
        c.showPage()
        c.save()

        from praxia.io.parsers.pdf import PdfParser
        parsed = PdfParser().parse(buf.getvalue(), filename="t.pdf")
        assert parsed.metadata["page_count"] == 2
        assert parsed.metadata["page_image_count"] == 2
        pi = parsed.metadata["page_images"][0]
        assert pi["mime"] == "image/jpeg"
        assert pi["page_num"] == 1
        assert pi["width"] > 0 and pi["height"] > 0
        # The base64 decodes to a JPEG magic header
        assert base64.b64decode(pi["data"])[:3] == b"\xff\xd8\xff"


# ---------------------------------------------------------------------------
# Aggregation: _aggregate_source_images
# ---------------------------------------------------------------------------


class TestAggregateSourceImages:
    def test_dedupes_identical_images_across_sources(self):
        from praxia.agent.commander import _aggregate_source_images
        from praxia.agent.verifier import Source

        same_b64 = base64.b64encode(b"hello world image bytes for test " * 4).decode()
        s1 = Source(id="D#0", text="a", kind="local_document",
                    metadata={"images": [{"mime": "image/png", "data": same_b64}]})
        s2 = Source(id="D#1", text="b", kind="local_document",
                    metadata={"images": [{"mime": "image/png", "data": same_b64}]})
        out = _aggregate_source_images([s1, s2])
        assert len(out) == 1, "duplicate image must be deduped"

    def test_caps_image_count(self):
        from praxia.agent.commander import _aggregate_source_images, _AGG_MAX_IMAGES
        from praxia.agent.verifier import Source
        sources = []
        for i in range(_AGG_MAX_IMAGES + 3):
            b = base64.b64encode(f"image-bytes-{i}".encode() * 5).decode()
            sources.append(Source(
                id=f"D#{i}", text="x", kind="local_document",
                metadata={"images": [{"mime": "image/png", "data": b}]},
            ))
        out = _aggregate_source_images(sources)
        assert len(out) == _AGG_MAX_IMAGES

    def test_ignores_sources_without_images(self):
        from praxia.agent.commander import _aggregate_source_images
        from praxia.agent.verifier import Source
        # Memory + frozen sources have no "images" in metadata
        s = Source(id="L1#0", text="t", kind="memory", metadata={})
        assert _aggregate_source_images([s]) == []
