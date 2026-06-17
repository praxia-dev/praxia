"""Tests for the vision-LLM PPTX review loop (Phase C-5).

These tests exercise the LLM-side logic only — we mock the vision
LLM and feed in pre-fabricated PNGs. The LibreOffice + pypdfium2
pipeline (``pptx_to_png.render_pptx_to_pngs``) is integration-
tested elsewhere because it requires an external binary.
"""
from __future__ import annotations

import json
from pathlib import Path
from typing import Any

import pytest


PIL = pytest.importorskip("PIL")  # tests need Pillow to fab PNGs

from praxia.io.document_reviewer import (
    DeckReview,
    PptxReviewer,
    SlideReview,
    _safe_parse_json,
)


# ---------------------------------------------------------------------------
# Test doubles
# ---------------------------------------------------------------------------


class _FakeVisionLLM:
    """Stub LLM that returns scripted JSON for each .complete() call.

    Pass either a single string (used for every call) or a list (one
    response per call, consumed in order). When the script runs out
    we raise to simulate a provider error so the SlideReview.error
    path is exercised.
    """

    def __init__(self, *, text: str | list[str] | None = None,
                 raise_after: int | None = None,
                 raise_exc: Exception | None = None) -> None:
        if isinstance(text, str):
            self._queue = [text]
            self._single = True
        else:
            self._queue = list(text or [])
            self._single = False
        self._raise_after = raise_after
        self._raise_exc = raise_exc or RuntimeError("vision LLM exhausted")
        self.calls: list[dict[str, Any]] = []

    def complete(self, messages, *, max_tokens=None, response_format=None,
                 temperature=None, **_kw):
        self.calls.append({
            "messages": messages,
            "max_tokens": max_tokens,
            "response_format": response_format,
            "temperature": temperature,
        })
        if self._raise_after is not None and len(self.calls) > self._raise_after:
            raise self._raise_exc
        if self._single:
            text = self._queue[0]
        elif self._queue:
            text = self._queue.pop(0)
        else:
            raise self._raise_exc
        return type("Resp", (), {"text": text})()


def _make_dummy_png(path: Path, *, color=(31, 58, 138)) -> None:
    """Write a tiny 1280x720 PNG so the reviewer has something to
    base64-encode."""
    from PIL import Image
    Image.new("RGB", (1280, 720), color=color).save(path, "PNG")


# ---------------------------------------------------------------------------
# JSON parser tolerance
# ---------------------------------------------------------------------------


class TestSafeParseJson:
    def test_bare_json_object(self):
        assert _safe_parse_json('{"a": 1}') == {"a": 1}

    def test_strips_code_fence(self):
        wrapped = '```json\n{"slide_score": 8}\n```'
        assert _safe_parse_json(wrapped) == {"slide_score": 8}

    def test_extracts_first_object_from_chatter(self):
        msg = 'Sure! Here is the review: {"slide_score": 6, "issues": []}. Hope that helps.'
        assert _safe_parse_json(msg) == {"slide_score": 6, "issues": []}

    def test_empty_string_returns_empty_dict(self):
        assert _safe_parse_json("") == {}

    def test_garbage_returns_empty_dict(self):
        assert _safe_parse_json("not json at all") == {}


# ---------------------------------------------------------------------------
# Single-slide review
# ---------------------------------------------------------------------------


class TestPptxReviewerSingleSlide:
    def test_clean_json_produces_full_slide_review(self, tmp_path: Path):
        png = tmp_path / "slide-01.png"
        _make_dummy_png(png)
        good_review = json.dumps({
            "slide_score": 8.5,
            "dimension_scores": {
                "typography": 9, "palette": 8, "layout": 8,
                "density": 9, "hierarchy": 8,
            },
            "issues": [
                {"category": "density", "severity": "low",
                 "description": "Bullets a touch long",
                 "suggested_fix": "Trim bullets 3 and 4 to one line each."}
            ],
            "strengths": ["Clean indigo + amber palette",
                          "Sentence-form headline"],
        })
        reviewer = PptxReviewer(llm=_FakeVisionLLM(text=good_review))
        deck = reviewer.review_pngs([png])
        assert isinstance(deck, DeckReview)
        assert len(deck.slides) == 1
        sr = deck.slides[0]
        assert sr.slide_number == 1
        assert sr.score == pytest.approx(8.5)
        assert sr.dimension_scores["typography"] == pytest.approx(9.0)
        assert len(sr.issues) == 1
        assert sr.issues[0]["severity"] == "low"
        assert "indigo" in sr.strengths[0].lower()

    def test_llm_error_captured_in_slide_review(self, tmp_path: Path):
        png = tmp_path / "slide-01.png"
        _make_dummy_png(png)
        reviewer = PptxReviewer(llm=_FakeVisionLLM(
            text=[], raise_after=0,
            raise_exc=ConnectionError("provider down"),
        ))
        deck = reviewer.review_pngs([png])
        sr = deck.slides[0]
        assert sr.score == 0.0
        assert "provider down" in sr.error

    def test_garbage_json_yields_zero_score(self, tmp_path: Path):
        png = tmp_path / "slide-01.png"
        _make_dummy_png(png)
        reviewer = PptxReviewer(llm=_FakeVisionLLM(text="not json at all"))
        deck = reviewer.review_pngs([png])
        sr = deck.slides[0]
        assert sr.score == 0.0
        assert sr.issues == []
        assert sr.error == ""  # didn't error, just produced no usable data

    def test_request_uses_temperature_zero_and_json_format(self, tmp_path: Path):
        png = tmp_path / "slide-01.png"
        _make_dummy_png(png)
        llm = _FakeVisionLLM(text='{"slide_score": 5}')
        reviewer = PptxReviewer(llm=llm)
        reviewer.review_pngs([png])
        assert llm.calls[0]["temperature"] == 0.0
        assert llm.calls[0]["response_format"] == "json"

    def test_image_attached_as_data_url(self, tmp_path: Path):
        png = tmp_path / "slide-01.png"
        _make_dummy_png(png)
        llm = _FakeVisionLLM(text='{"slide_score": 7}')
        reviewer = PptxReviewer(llm=llm)
        reviewer.review_pngs([png])
        # The user-role message should contain a multimodal block
        # with the PNG as a base64 data URL.
        msgs = llm.calls[0]["messages"]
        user_msg = next(m for m in msgs if m["role"] == "user")
        assert isinstance(user_msg["content"], list)
        types = [block.get("type") for block in user_msg["content"]]
        assert "image_url" in types
        image_block = next(b for b in user_msg["content"] if b["type"] == "image_url")
        assert image_block["image_url"]["url"].startswith("data:image/png;base64,")


# ---------------------------------------------------------------------------
# Deck-wide aggregation
# ---------------------------------------------------------------------------


class TestDeckReview:
    def test_overall_score_is_mean_across_slides(self, tmp_path: Path):
        pngs = []
        for i in range(3):
            p = tmp_path / f"slide-0{i+1}.png"
            _make_dummy_png(p)
            pngs.append(p)
        scripted = [
            json.dumps({"slide_score": s, "issues": []})
            for s in (6.0, 8.0, 10.0)
        ]
        deck = PptxReviewer(llm=_FakeVisionLLM(text=scripted)).review_pngs(pngs)
        assert deck.overall_score == pytest.approx(8.0)

    def test_high_severity_issues_aggregated_with_slide_numbers(self, tmp_path: Path):
        pngs = [tmp_path / f"slide-0{i+1}.png" for i in range(2)]
        for p in pngs:
            _make_dummy_png(p)
        scripted = [
            json.dumps({
                "slide_score": 5,
                "issues": [
                    {"category": "palette", "severity": "high",
                     "description": "Black on white",
                     "suggested_fix": "Use indigo title color"},
                    {"category": "density", "severity": "low",
                     "description": "Could trim", "suggested_fix": "—"},
                ],
            }),
            json.dumps({
                "slide_score": 9,
                "issues": [
                    {"category": "hierarchy", "severity": "high",
                     "description": "Title competes with body",
                     "suggested_fix": "Increase title to 36pt"},
                ],
            }),
        ]
        deck = PptxReviewer(llm=_FakeVisionLLM(text=scripted)).review_pngs(pngs)
        hi = deck.high_severity_issues
        assert len(hi) == 2
        # Slide numbers are injected so the agent can quote them
        assert {h["slide"] for h in hi} == {1, 2}

    def test_to_dict_round_trips_through_json(self, tmp_path: Path):
        png = tmp_path / "slide-01.png"
        _make_dummy_png(png)
        deck = PptxReviewer(llm=_FakeVisionLLM(text='{"slide_score": 7.5}')).review_pngs([png])
        d = deck.to_dict()
        # Should be JSON-serialisable end to end.
        roundtrip = json.loads(json.dumps(d))
        assert roundtrip["overall_score"] == pytest.approx(7.5)
        assert roundtrip["slides"][0]["score"] == pytest.approx(7.5)
        assert "summary" in roundtrip

    def test_empty_png_list_returns_zero_overall(self):
        deck = PptxReviewer(llm=_FakeVisionLLM(text="never reached")).review_pngs([])
        assert deck.overall_score == 0.0
        assert deck.slides == []


# ---------------------------------------------------------------------------
# Dependency probe
# ---------------------------------------------------------------------------


class TestDependencyProbe:
    def test_check_dependencies_returns_ok_or_reason(self):
        from praxia.io.pptx_to_png import check_dependencies
        ok, reason = check_dependencies()
        # On a typical CI box neither LibreOffice nor pypdfium2 are
        # guaranteed present; just assert the contract.
        assert isinstance(ok, bool)
        if not ok:
            assert isinstance(reason, str) and reason


# ---------------------------------------------------------------------------
# Text-mode reviewer — uses python-pptx, no LibreOffice required
# ---------------------------------------------------------------------------


pytest.importorskip("pptx")


def _build_minimal_pptx(path: Path) -> None:
    """Build a small styled deck via the slide_templates so the
    describer has something realistic to walk. No LibreOffice."""
    from pptx import Presentation
    from pptx.util import Inches
    from praxia.io.slide_templates import (
        cover_slide, kpi_slide, bullets_slide,
    )
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    cover_slide(prs, title="Q3 review", subtitle="Numbers", kicker="2026")
    kpi_slide(prs, title="Highlights", kpis=[
        {"label": "ARR", "value": "$412k", "delta": "+$32k"},
        {"label": "NPS", "value": "42", "delta": "+4"},
    ])
    bullets_slide(prs, title="Roadmap", bullets=["Ship X", "Test Y"])
    prs.save(str(path))


class TestTextModeReviewer:
    def test_describe_pptx_returns_one_string_per_slide(self, tmp_path: Path):
        from praxia.io.pptx_describer import describe_pptx
        deck = tmp_path / "deck.pptx"
        _build_minimal_pptx(deck)
        descs = describe_pptx(deck)
        assert len(descs) == 3
        # Each description carries the structural fields the LLM
        # needs: a slide marker, layout name, slide size, shape blocks
        for d in descs:
            assert d.startswith("=== Slide")
            assert "slide_size:" in d
            assert "Shape 1" in d

    def test_describer_captures_palette_hex_codes(self, tmp_path: Path):
        from praxia.io.pptx_describer import describe_pptx
        deck = tmp_path / "deck.pptx"
        _build_minimal_pptx(deck)
        descs = describe_pptx(deck)
        # Cover slide has the indigo top panel
        assert "#1F3A8A" in descs[0]
        # KPI slide has the amber accent bar
        assert "#F59E0B" in descs[1]

    def test_text_mode_review_via_explicit_mode(self, tmp_path: Path):
        from praxia.io.document_reviewer import PptxReviewer
        deck_path = tmp_path / "deck.pptx"
        _build_minimal_pptx(deck_path)
        scripted = [
            json.dumps({
                "slide_score": 8 - i,
                "dimension_scores": {"typography": 8, "palette": 9},
                "issues": [],
                "strengths": ["clean palette"],
            }) for i in range(3)
        ]
        llm = _FakeVisionLLM(text=scripted)
        reviewer = PptxReviewer(llm=llm, mode="text")
        deck = reviewer.review(deck_path)
        assert deck.mode == "text"
        assert len(deck.slides) == 3
        # All three slides should have non-zero scores
        assert all(sr.score > 0 for sr in deck.slides)
        # Verify the user-role message embeds the structural description
        # (it should contain the "=== Slide N ===" header)
        user_msg = next(m for m in llm.calls[0]["messages"] if m["role"] == "user")
        # Text mode uses a plain string content (no multimodal)
        assert isinstance(user_msg["content"], str)
        assert "=== Slide" in user_msg["content"]

    def test_auto_mode_falls_back_to_text_when_vision_unavailable(
        self, tmp_path: Path, monkeypatch
    ):
        """Auto mode should pick text mode when the LibreOffice+pypdfium2
        path isn't available, and stamp the result with an upgrade hint."""
        from praxia.io import document_reviewer as dr
        from praxia.io.document_reviewer import PptxReviewer
        deck_path = tmp_path / "deck.pptx"
        _build_minimal_pptx(deck_path)

        # Force vision-deps probe to report unavailable.
        monkeypatch.setattr(
            "praxia.io.pptx_to_png.check_dependencies",
            lambda: (False, "(mocked: LibreOffice missing)"),
        )
        llm = _FakeVisionLLM(text=json.dumps({"slide_score": 7}))
        reviewer = PptxReviewer(llm=llm, mode="auto")
        deck = reviewer.review(deck_path)
        assert deck.mode == "text"
        assert deck.upgrade_hint
        assert "LibreOffice" in deck.upgrade_hint or "vision" in deck.upgrade_hint.lower()

    def test_to_dict_carries_mode_and_upgrade_hint(self, tmp_path: Path, monkeypatch):
        from praxia.io.document_reviewer import PptxReviewer
        deck_path = tmp_path / "deck.pptx"
        _build_minimal_pptx(deck_path)
        monkeypatch.setattr(
            "praxia.io.pptx_to_png.check_dependencies",
            lambda: (False, "(mocked)"),
        )
        llm = _FakeVisionLLM(text=json.dumps({"slide_score": 7}))
        reviewer = PptxReviewer(llm=llm, mode="auto")
        result = reviewer.review(deck_path).to_dict()
        assert result["mode"] == "text"
        assert "upgrade_hint" in result

    def test_invalid_mode_rejected(self):
        from praxia.io.document_reviewer import PptxReviewer
        with pytest.raises(ValueError):
            PptxReviewer(llm=_FakeVisionLLM(text=""), mode="bogus")  # type: ignore[arg-type]
