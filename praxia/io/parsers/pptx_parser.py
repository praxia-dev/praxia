"""Microsoft PowerPoint (.pptx) parser — uses python-pptx for text +
tables and a bare zipfile read for the embedded `ppt/media/*` images
so the vision LLM can see charts / diagrams / screenshots that the
.pptx contains."""
from __future__ import annotations

import base64
import io
import re
import zipfile
from typing import Any

from praxia.io.parsers.base import ParsedFile

# Same caps as the .docx parser for consistency.
_MAX_PER_IMAGE_BYTES = 5 * 1024 * 1024
_MAX_TOTAL_IMAGE_BYTES = 20 * 1024 * 1024
_MEDIA_PATH_RE = re.compile(r"^ppt/media/.*\.(png|jpe?g|gif|webp)$", re.IGNORECASE)


def _ext_to_mime(name: str) -> str | None:
    n = name.lower()
    if n.endswith(".png"):  return "image/png"
    if n.endswith(".jpg") or n.endswith(".jpeg"): return "image/jpeg"
    if n.endswith(".gif"):  return "image/gif"
    if n.endswith(".webp"): return "image/webp"
    return None


def _extract_embedded_images(data: bytes) -> list[dict[str, Any]]:
    """Return a list of {name, mime, data (base64), bytes_size} from
    the .pptx's `ppt/media/` folder. Same caps + behaviour as the
    DOCX parser's helper — see that module for the rationale."""
    out: list[dict[str, Any]] = []
    total = 0
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            for info in zf.infolist():
                if not _MEDIA_PATH_RE.match(info.filename):
                    continue
                if info.file_size > _MAX_PER_IMAGE_BYTES:
                    continue
                if total + info.file_size > _MAX_TOTAL_IMAGE_BYTES:
                    break
                mime = _ext_to_mime(info.filename)
                if not mime:
                    continue
                try:
                    raw = zf.read(info.filename)
                except KeyError:
                    continue
                out.append({
                    "name": info.filename.split("/")[-1],
                    "mime": mime,
                    "data": base64.b64encode(raw).decode("ascii"),
                    "bytes_size": info.file_size,
                })
                total += info.file_size
    except zipfile.BadZipFile:
        return []
    return out


class PptxParser:
    name = "pptx"

    def parse(self, data: bytes, *, filename: str, **kwargs: Any) -> ParsedFile:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError(
                "PowerPoint (.pptx) parsing requires `python-pptx`. Install with:\n"
                '  pip install "praxia[office]"'
            ) from e

        prs = Presentation(io.BytesIO(data))
        sections: list[tuple[str, str]] = []
        full_text: list[str] = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = []
            slide_title: str | None = None

            for shape in slide.shapes:
                # Title
                if shape.has_text_frame:
                    tf = shape.text_frame
                    text = tf.text.strip()
                    if not text:
                        continue
                    if (
                        slide_title is None
                        and getattr(shape, "is_placeholder", False)
                        and "title" in str(shape.placeholder_format.type or "").lower()
                    ):
                        slide_title = text
                    slide_lines.append(text)
                # Tables in slides
                if shape.has_table:
                    rows: list[list[str]] = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        slide_lines.append(
                            "| " + " | ".join(rows[0]) + " |"
                        )
                        slide_lines.append(
                            "|" + "|".join(["---"] * len(rows[0])) + "|"
                        )
                        for row in rows[1:]:
                            slide_lines.append("| " + " | ".join(row) + " |")
                # Notes
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            heading = slide_title or f"Slide {i}"
            section_body = "\n".join(slide_lines)
            if notes_text:
                section_body += f"\n\n[Speaker notes]\n{notes_text}"
            sections.append((heading, section_body))

            full_text.append(f"\n## Slide {i}: {heading}\n{section_body}")

        # Embedded images — alpha20 multi-modal pipeline. Each slide's
        # charts / pictures / diagrams that were inserted as images
        # land here.
        embedded = _extract_embedded_images(data)

        return ParsedFile(
            filename=filename,
            content="\n".join(full_text),
            metadata={
                "slide_count": len(prs.slides),
                "embedded_images": embedded,
                "embedded_image_count": len(embedded),
            },
            sections=sections,
        )
