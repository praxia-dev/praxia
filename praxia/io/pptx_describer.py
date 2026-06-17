"""Text-only structural description of a .pptx, used by the
text-mode fallback in :class:`praxia.io.document_reviewer.PptxReviewer`.

Why it exists:
    The vision-mode reviewer needs LibreOffice + pypdfium2 to
    rasterise each slide before a vision LLM sees it. Both are
    optional installs. For users who don't want to install 400 MB
    of LibreOffice, we instead walk the python-pptx object model
    and emit a structured text description per slide, then ask a
    *text* LLM to critique that. The text path can't see actual
    rendering issues (overlap, contrast on real backgrounds, font
    fallbacks) but covers typography / palette / density / rough
    layout reasoning at ~70-80% of vision-mode quality.

Output shape (one string per slide):

    === Slide N ===
    layout: <slide_layout.name>
    slide_size: 13.3 x 7.5 in

      Shape 1 (textbox):
        position: (0.6, 0.5) in, size: 12.0 x 0.8 in
        fill: #1F3A8A
        text:
          [0.0] Segoe UI Semibold 36pt #FFFFFF bold align=left "Q3 Sales Review"
          [0.1] Segoe UI 18pt #FFFFFF align=left "ARR progression"

      Shape 2 (picture):
        position: (0.6, 1.8) in, size: 8.0 x 4.8 in
        kind: image (likely a chart)

      Shape 3 (rectangle):
        position: (8.9, 1.8) in, size: 3.8 x 4.8 in
        fill: #1F3A8A

A text LLM with the critique rubric can score this without seeing
the actual PNG.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

_log = logging.getLogger(__name__)


_EMU_PER_INCH = 914400


def describe_pptx(pptx_path: Path) -> list[str]:
    """Return one text description per slide in the deck.

    Slides are emitted in order. Each string is self-contained so
    callers can review one slide at a time.

    Raises:
        ImportError: if python-pptx is not installed. (Unlikely in
        practice — we just rendered the deck with it.)
    """
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError(
            "python-pptx is required for text-mode review. "
            "Install with: pip install 'praxia[office]'"
        ) from e

    prs = Presentation(str(pptx_path))
    width_in = prs.slide_width / _EMU_PER_INCH if prs.slide_width else 0.0
    height_in = prs.slide_height / _EMU_PER_INCH if prs.slide_height else 0.0
    out: list[str] = []
    for i, slide in enumerate(prs.slides):
        lines = [f"=== Slide {i + 1} ==="]
        layout_name = getattr(getattr(slide, "slide_layout", None), "name", "") or ""
        if layout_name:
            lines.append(f"layout: {layout_name}")
        lines.append(f"slide_size: {width_in:.1f} x {height_in:.1f} in")
        bg = _slide_background_color(slide)
        if bg:
            lines.append(f"background: #{bg}")
        lines.append("")
        for j, shape in enumerate(slide.shapes):
            lines.append(_describe_shape(shape, j + 1))
        out.append("\n".join(lines))
    return out


def _describe_shape(shape, n: int) -> str:
    """Per-shape block. Captures kind, position, fill, and text runs
    (with font + colour + bold) — all the fields a critique LLM needs
    to reason about typography/palette/density."""
    kind = _shape_kind(shape)
    parts = [f"  Shape {n} ({kind}):"]

    pos = _shape_position(shape)
    if pos is not None:
        x, y, w, h = pos
        parts.append(f"    position: ({x:.1f}, {y:.1f}) in, size: {w:.1f} x {h:.1f} in")

    fill = _shape_fill_color(shape)
    if fill:
        parts.append(f"    fill: #{fill}")

    line_color = _shape_line_color(shape)
    if line_color:
        parts.append(f"    border: #{line_color}")

    if kind == "picture":
        # We can't read the embedded PNG content, but flag it as an
        # image so the reviewer knows there's a visual element here.
        parts.append("    kind: image (possibly chart / diagram / photo)")
    elif _has_text(shape):
        parts.append("    text:")
        for p_idx, para in enumerate(shape.text_frame.paragraphs):
            align = _safe_alignment_name(para)
            for r_idx, run in enumerate(para.runs):
                font_descr = _describe_font(run.font)
                preview = (run.text or "").replace("\n", " ").strip()
                if not preview:
                    continue
                if len(preview) > 100:
                    preview = preview[:97] + "…"
                align_part = f" align={align}" if align else ""
                parts.append(
                    f"      [{p_idx}.{r_idx}] {font_descr}{align_part} \"{preview}\""
                )
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Internal extractors — defensive against python-pptx returning None for
# inherited / not-yet-resolved values.
# ---------------------------------------------------------------------------


def _shape_kind(shape) -> str:
    """Short human-readable kind label."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return "shape"
    st = getattr(shape, "shape_type", None)
    if st is None:
        return "shape"
    mapping = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "autoshape",
        MSO_SHAPE_TYPE.TEXT_BOX: "textbox",
        MSO_SHAPE_TYPE.PICTURE: "picture",
        MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
        MSO_SHAPE_TYPE.GROUP: "group",
        MSO_SHAPE_TYPE.LINE: "line",
        MSO_SHAPE_TYPE.FREEFORM: "freeform",
        MSO_SHAPE_TYPE.TABLE: "table",
        MSO_SHAPE_TYPE.CHART: "chart",
        MSO_SHAPE_TYPE.MEDIA: "media",
    }
    return mapping.get(st, str(st).lower().replace("mso_shape_type.", ""))


def _shape_position(shape) -> tuple[float, float, float, float] | None:
    """(left, top, width, height) in inches, or None if any are
    inherited / unset."""
    try:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
    except (ValueError, AttributeError):
        return None
    if None in (left, top, width, height):
        return None
    return (
        float(left) / _EMU_PER_INCH,
        float(top) / _EMU_PER_INCH,
        float(width) / _EMU_PER_INCH,
        float(height) / _EMU_PER_INCH,
    )


def _safe_rgb_hex(fill_proxy) -> str | None:
    """Extract a solid-fill RGB as hex. Returns None for None / picture /
    pattern / theme-coloured / inherited / no-fill cases — python-pptx
    raises TypeError when ``fore_color`` is accessed on a non-solid fill.
    """
    if fill_proxy is None:
        return None
    try:
        from pptx.dml.color import RGBColor
        from pptx.enum.dml import MSO_FILL_TYPE
    except ImportError:
        return None
    # `.type` is None until the fill descriptor is resolved; calling
    # `.fore_color` on a non-solid fill raises TypeError. Probe both.
    try:
        fill_type = fill_proxy.type
    except Exception:
        fill_type = None
    if fill_type is not None and fill_type != MSO_FILL_TYPE.SOLID:
        return None
    try:
        fore = fill_proxy.fore_color
        rgb = fore.rgb
    except (TypeError, AttributeError):
        return None
    if rgb is None or not isinstance(rgb, RGBColor):
        return None
    try:
        # python-pptx's RGBColor is a `str` subclass holding the 6-hex
        # representation already — int() coercion fails on it. We
        # uppercase for consistency.
        return str(rgb).upper()
    except Exception:
        return None


def _shape_fill_color(shape) -> str | None:
    """Hex (no #) of the solid fill, or None for non-solid /
    inherited / theme-coloured fills (which we can't trivially resolve
    without the theme XML)."""
    return _safe_rgb_hex(getattr(shape, "fill", None))


def _shape_line_color(shape) -> str | None:
    """Border colour (hex without #). None when not solid / inherited."""
    line = getattr(shape, "line", None)
    if line is None:
        return None
    try:
        rgb = getattr(getattr(line, "color", None), "rgb", None)
    except (TypeError, AttributeError):
        return None
    if rgb is None:
        return None
    try:
        # python-pptx's RGBColor is a `str` subclass holding the 6-hex
        # representation already — int() coercion fails on it. We
        # uppercase for consistency.
        return str(rgb).upper()
    except Exception:
        return None


def _slide_background_color(slide) -> str | None:
    bg = getattr(slide, "background", None)
    return _safe_rgb_hex(getattr(bg, "fill", None)) if bg is not None else None


def _has_text(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    return bool((shape.text_frame.text or "").strip())


def _describe_font(font) -> str:
    """Compact 'Family Size#hex (bold/italic)' representation."""
    name = getattr(font, "name", None) or "(default)"
    size_obj = getattr(font, "size", None)
    if size_obj is not None:
        try:
            size = f"{int(size_obj.pt)}pt"
        except Exception:
            size = ""
    else:
        size = ""
    color = ""
    color_obj = getattr(font, "color", None)
    if color_obj is not None:
        try:
            rgb = color_obj.rgb
        except (TypeError, AttributeError):
            rgb = None
        if rgb is not None:
            try:
                color = f" #{str(rgb).upper()}"
            except Exception:
                color = ""
    modifiers = []
    if getattr(font, "bold", None):
        modifiers.append("bold")
    if getattr(font, "italic", None):
        modifiers.append("italic")
    mod_part = f" {'+'.join(modifiers)}" if modifiers else ""
    return f"{name} {size}{color}{mod_part}".strip()


def _safe_alignment_name(para) -> str | None:
    """Convert PP_ALIGN to a string the LLM can read. Returns None if
    the paragraph inherits its alignment."""
    align = getattr(para, "alignment", None)
    if align is None:
        return None
    name = getattr(align, "name", None)
    if name:
        return name.lower()
    return str(align).lower().replace("pp_align.", "").split()[0]
