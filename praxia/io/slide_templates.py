"""High-level slide builders that produce a styled deck.

The :class:`PptxExporter` falls back to these when the input is a
structured spec (see ``_iter_spec_slides``); the LLM is also free to
feed structured specs via the new ``slides`` Markdown extension —
fences of the form ``\`\`\`slide:cover`` / ``slide:kpi`` / etc.

What each helper produces:

    cover_slide      Big H1, optional subtitle and accent kicker bar.
                     One per deck; goes first.

    section_slide    Indigo-on-white wedge with a bold label. Use to
                     break the deck into chapters when there are
                     more than ~6 content slides.

    kpi_slide        Up to 4 KPI tiles in a row. Each tile is a
                     big number + label + delta. Use this instead
                     of bullet-listing percentages.

    chart_slide      Title + matplotlib chart (saved as PNG and
                     embedded) + 2-3 "takeaway" bullets on the right.

    bullets_slide    Title + bulleted body. The fallback when the
                     content doesn't deserve a chart or KPI.

Rationale:
    Hand-built `tmp/build_q3_sales_review.py` is the gold standard
    for what a Praxia deck should look like. Capturing those patterns
    here gives the exporter (and any LLM that writes structured
    slide specs) a consistent set of building blocks. This is the
    "lego" layer; :mod:`slide_style` is the colour/font/typography
    layer.

These helpers require ``python-pptx``. Charts additionally require
``matplotlib`` — if matplotlib is missing, ``chart_slide`` falls
back to a styled table rendering of the data.
"""
from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any, Iterable, Sequence

from praxia.io.slide_style import (
    DEFAULT_STYLE,
    SlideStyle,
    font_for,
)


# ---------------------------------------------------------------------------
# Internal helpers — every public template uses these so the visual
# language stays consistent across cover / section / kpi / chart / bullets.
# ---------------------------------------------------------------------------


def _rgb(triplet):
    """Convert (r,g,b) tuple to pptx's RGBColor object lazily."""
    from pptx.dml.color import RGBColor
    return RGBColor(*triplet)


def _pt(size: int):
    from pptx.util import Pt
    return Pt(size)


def _inches(value):
    from pptx.util import Inches
    return Inches(value)


def _blank_slide(prs, *, fill: tuple[int, int, int] | None = None):
    """Create a blank slide and paint a solid background.

    We always use the blank layout (index 6 on the default template)
    because the built-in title/content layouts pull from a master
    that uses Calibri 28pt — we don't want that polluting our
    typography.
    """
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    if fill is not None:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(fill)
    return slide


def _add_text(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    bold: bool = False,
    color: tuple[int, int, int] = (17, 24, 39),
    align: str = "left",
):
    """Add a textbox with consistent formatting. ``align`` is one of
    ``"left" / "center" / "right"``."""
    from pptx.enum.text import PP_ALIGN
    box = slide.shapes.add_textbox(_inches(left), _inches(top),
                                   _inches(width), _inches(height))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size = _pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    reg, bold_name = font_for(text)
    run.font.name = bold_name if bold else reg
    return box


def _add_accent_bar(slide, *, left: float, top: float, width: float, color):
    """Thin accent rule (3 pt tall) used under cover/section titles."""
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        _inches(left), _inches(top), _inches(width), _inches(0.05),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(color)
    bar.line.fill.background()
    return bar


def _add_panel(slide, *, left, top, width, height, fill, border=None):
    """Solid filled rectangle used as a card / panel background."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        _inches(left), _inches(top), _inches(width), _inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if border:
        shape.line.color.rgb = _rgb(border)
        shape.line.width = _pt(1)
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Public templates
# ---------------------------------------------------------------------------


def cover_slide(
    prs,
    *,
    title: str,
    subtitle: str | None = None,
    kicker: str | None = None,
    style: SlideStyle = DEFAULT_STYLE,
):
    """Title slide: indigo background, white title, optional subtitle.

    Layout:
        | (top half — indigo background)              |
        |   KICKER (amber, small)                     |
        |   TITLE  (white, very large)                |
        |   ───── (amber accent bar)                  |
        | (bottom half — white)                       |
        |   SUBTITLE  (grey-700)                      |

    Use exactly one per deck, as slide 1.
    """
    slide = _blank_slide(prs)

    # Top half: indigo block
    _add_panel(slide, left=0, top=0, width=13.33, height=4.5,
               fill=style.primary)

    if kicker:
        _add_text(slide, kicker.upper(),
                  left=style.margin_left, top=0.6, width=12.0, height=0.4,
                  size=14, bold=True, color=style.accent, align="left")

    _add_text(slide, title,
              left=style.margin_left, top=1.4, width=12.0, height=2.4,
              size=style.cover_title_pt, bold=True,
              color=(255, 255, 255), align="left")

    _add_accent_bar(slide, left=style.margin_left, top=3.9,
                    width=1.6, color=style.accent)

    if subtitle:
        _add_text(slide, subtitle,
                  left=style.margin_left, top=5.0, width=12.0, height=1.4,
                  size=20, color=style.text, align="left")

    return slide


def section_slide(
    prs,
    *,
    label: str,
    section_number: int | None = None,
    style: SlideStyle = DEFAULT_STYLE,
):
    """Section divider. White background, amber wedge on the left, big
    indigo label. Use between chapters when the deck is >6 slides."""
    slide = _blank_slide(prs, fill=style.background)
    _add_panel(slide, left=0, top=0, width=0.8, height=7.5,
               fill=style.accent)

    if section_number is not None:
        _add_text(slide, f"{section_number:02d}",
                  left=1.4, top=2.4, width=2.0, height=0.6,
                  size=18, bold=True, color=style.accent, align="left")

    _add_text(slide, label,
              left=1.4, top=3.0, width=11.0, height=1.5,
              size=style.section_pt, bold=True,
              color=style.primary, align="left")
    _add_accent_bar(slide, left=1.4, top=4.5, width=1.6, color=style.accent)
    return slide


def kpi_slide(
    prs,
    *,
    title: str,
    kpis: Sequence[dict[str, Any]],
    style: SlideStyle = DEFAULT_STYLE,
):
    """KPI row: up to 4 large numeric tiles.

    Each KPI is a dict with:
        label   short caption (e.g. "ARR", "NPS")
        value   the big number as a string (e.g. "$412k", "42")
        delta   optional delta string (e.g. "+$32k", "+4 pts");
                positive→green, negative→red, "—"→muted
    """
    if not kpis:
        raise ValueError("kpi_slide requires at least one KPI")
    slide = _blank_slide(prs, fill=style.background)

    _add_text(slide, title,
              left=style.margin_left, top=style.margin_top,
              width=12.0, height=0.8,
              size=style.title_pt, bold=True, color=style.primary)
    _add_accent_bar(slide, left=style.margin_left, top=1.35,
                    width=1.2, color=style.accent)

    n = min(len(kpis), 4)
    total_w = 13.33 - 2 * style.margin_left
    gutter = 0.25
    tile_w = (total_w - gutter * (n - 1)) / n
    tile_h = 3.0
    tile_top = 2.3
    for i, kpi in enumerate(kpis[:n]):
        x = style.margin_left + i * (tile_w + gutter)
        _add_panel(slide, left=x, top=tile_top, width=tile_w, height=tile_h,
                   fill=style.panel_bg)
        _add_text(slide, str(kpi.get("label", "")),
                  left=x + 0.2, top=tile_top + 0.2,
                  width=tile_w - 0.4, height=0.4,
                  size=style.kpi_label_pt, bold=True,
                  color=style.text_muted)
        _add_text(slide, str(kpi.get("value", "")),
                  left=x + 0.2, top=tile_top + 0.7,
                  width=tile_w - 0.4, height=1.4,
                  size=style.kpi_value_pt, bold=True,
                  color=style.primary)
        delta = kpi.get("delta")
        if delta:
            s = str(delta).strip()
            if s.startswith("-"):
                dcolor = (220, 38, 38)   # red
            elif s.startswith("+"):
                dcolor = (16, 185, 129)  # green
            else:
                dcolor = style.text_muted
            _add_text(slide, s,
                      left=x + 0.2, top=tile_top + 2.2,
                      width=tile_w - 0.4, height=0.5,
                      size=style.body_pt, bold=True, color=dcolor)
    return slide


def chart_slide(
    prs,
    *,
    title: str,
    chart_type: str,
    labels: Sequence[str],
    values: Sequence[Sequence[float]] | Sequence[float],
    series_names: Sequence[str] | None = None,
    takeaways: Sequence[str] = (),
    y_label: str | None = None,
    style: SlideStyle = DEFAULT_STYLE,
):
    """Title + matplotlib chart on the left + 2-3 takeaway bullets on
    the right.

    ``chart_type``: ``"bar"`` | ``"line"`` | ``"hbar"`` (horizontal bar)
    ``values``: single series → flat list; multi-series → list of lists
    """
    slide = _blank_slide(prs, fill=style.background)

    _add_text(slide, title,
              left=style.margin_left, top=style.margin_top,
              width=12.0, height=0.8,
              size=style.title_pt, bold=True, color=style.primary)
    _add_accent_bar(slide, left=style.margin_left, top=1.35,
                    width=1.2, color=style.accent)

    chart_png = _render_chart_png(
        chart_type=chart_type,
        labels=labels,
        values=values,
        series_names=series_names,
        y_label=y_label,
        style=style,
    )
    chart_left = style.margin_left
    chart_top = 1.8
    chart_w = 8.0
    chart_h = 4.8
    if chart_png is not None:
        slide.shapes.add_picture(
            BytesIO(chart_png),
            _inches(chart_left), _inches(chart_top),
            width=_inches(chart_w), height=_inches(chart_h),
        )
    else:
        # matplotlib not installed; render a tabular fallback so the
        # slide isn't blank — the data still gets across.
        _render_table_fallback(
            slide, chart_left, chart_top, chart_w, chart_h,
            labels, values, series_names, style,
        )

    if takeaways:
        panel_left = chart_left + chart_w + 0.3
        panel_w = 13.33 - panel_left - style.margin_right
        _add_panel(slide, left=panel_left, top=chart_top,
                   width=panel_w, height=chart_h, fill=style.primary)
        _add_text(slide, "Takeaways",
                  left=panel_left + 0.25, top=chart_top + 0.2,
                  width=panel_w - 0.5, height=0.4,
                  size=16, bold=True, color=style.accent)
        body_top = chart_top + 0.9
        per_line = (chart_h - 1.1) / max(1, len(takeaways))
        for i, line in enumerate(takeaways):
            _add_text(slide, f"·  {line}",
                      left=panel_left + 0.25,
                      top=body_top + i * per_line,
                      width=panel_w - 0.5, height=per_line,
                      size=14, color=(255, 255, 255))
    return slide


def bullets_slide(
    prs,
    *,
    title: str,
    bullets: Sequence[str],
    style: SlideStyle = DEFAULT_STYLE,
):
    """Title + indented bullets. The fallback when the content is
    qualitative and doesn't justify a chart or KPI row."""
    slide = _blank_slide(prs, fill=style.background)

    _add_text(slide, title,
              left=style.margin_left, top=style.margin_top,
              width=12.0, height=0.8,
              size=style.title_pt, bold=True, color=style.primary)
    _add_accent_bar(slide, left=style.margin_left, top=1.35,
                    width=1.2, color=style.accent)

    if not bullets:
        return slide

    from pptx.util import Inches
    body = slide.shapes.add_textbox(
        Inches(style.margin_left), Inches(2.0),
        Inches(13.33 - 2 * style.margin_left), Inches(4.8),
    )
    tf = body.text_frame
    tf.word_wrap = True
    from pptx.enum.text import PP_ALIGN
    for i, item in enumerate(bullets[:7]):  # cap at 7; more = redesign
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"·  {item}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = _pt(style.body_pt)
        run.font.color.rgb = _rgb(style.text)
        reg, _ = font_for(item)
        run.font.name = reg
        p.line_spacing = style.body_line_spacing
    return slide


# ---------------------------------------------------------------------------
# Chart rendering — matplotlib is optional; if missing, fall back to
# a tabular slide that at least preserves the data.
# ---------------------------------------------------------------------------


def _render_chart_png(
    *,
    chart_type: str,
    labels: Sequence[str],
    values: Sequence[Sequence[float]] | Sequence[float],
    series_names: Sequence[str] | None,
    y_label: str | None,
    style: SlideStyle,
) -> bytes | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None
    # Pick a font that has CJK glyphs if any label / series name does.
    sample = " ".join(list(labels) + list(series_names or []))
    matplotlib.rcParams["font.family"] = (
        ["Yu Gothic", "Meiryo", "MS Gothic", "DejaVu Sans"]
        if any("぀" <= ch <= "鿿" for ch in sample) else
        ["Segoe UI", "Inter", "DejaVu Sans"]
    )

    # Normalise values into a 2D shape: rows = series, cols = labels.
    if values and isinstance(values[0], (int, float)):
        matrix = [list(values)]
        names = list(series_names or [""])
    else:
        matrix = [list(s) for s in values]
        names = list(series_names or [f"Series {i+1}" for i in range(len(matrix))])

    fig_w, fig_h = 9.6, 5.4
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=150)
    palette = [tuple(c / 255 for c in rgb) for rgb in style.chart_series]

    if chart_type == "line":
        for i, row in enumerate(matrix):
            ax.plot(labels, row, marker="o", linewidth=2.5,
                    color=palette[i % len(palette)],
                    label=names[i] if names[i] else None)
    elif chart_type == "hbar":
        import numpy as np
        idx = list(range(len(labels)))
        bar_h = 0.8 / max(1, len(matrix))
        for i, row in enumerate(matrix):
            offsets = [j + (i - (len(matrix) - 1) / 2) * bar_h for j in idx]
            ax.barh(offsets, row, height=bar_h,
                    color=palette[i % len(palette)],
                    label=names[i] if names[i] else None)
        ax.set_yticks(idx)
        ax.set_yticklabels(labels)
        ax.invert_yaxis()
    else:  # "bar" (default)
        import numpy as np
        idx = list(range(len(labels)))
        bar_w = 0.8 / max(1, len(matrix))
        for i, row in enumerate(matrix):
            offsets = [j + (i - (len(matrix) - 1) / 2) * bar_w for j in idx]
            ax.bar(offsets, row, width=bar_w,
                   color=palette[i % len(palette)],
                   label=names[i] if names[i] else None)
        ax.set_xticks(idx)
        ax.set_xticklabels(labels)

    if y_label:
        ax.set_ylabel(y_label, fontsize=11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(labelsize=11)
    if len(matrix) > 1 and any(names):
        ax.legend(frameon=False, fontsize=10, loc="best")

    fig.tight_layout(pad=1.5)
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


def _render_table_fallback(slide, left, top, width, height,
                           labels, values, series_names, style):
    """Last resort when matplotlib is missing. Lay out the data as a
    small table so the numbers at least make it onto the slide."""
    if values and isinstance(values[0], (int, float)):
        rows = [("", *labels), ("", *(str(v) for v in values))]
    else:
        rows = [("", *labels)]
        for name, row in zip(series_names or [], values):
            rows.append((str(name), *(str(v) for v in row)))
    n_rows = len(rows)
    n_cols = len(rows[0])
    cell_w = width / n_cols
    cell_h = height / max(1, n_rows)
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            color = style.primary if r == 0 else style.text
            _add_text(slide, str(cell),
                      left=left + c * cell_w, top=top + r * cell_h,
                      width=cell_w, height=cell_h,
                      size=14, bold=(r == 0), color=color)
